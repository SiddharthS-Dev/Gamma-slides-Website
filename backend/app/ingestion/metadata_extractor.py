"""Metadata extractor — extracts title, author, slide count from PPTX and PDF files."""

import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


def extract_metadata(file_path: Path) -> dict:
    """Extract metadata from a presentation file.

    Returns a dict with keys: title, author, description, slide_count, created_date, modified_date
    """
    ext = file_path.suffix.lower()

    metadata = {}
    if ext == ".pptx":
        metadata = _extract_pptx_metadata(file_path)
    elif ext == ".pdf":
        metadata = _extract_pdf_metadata(file_path)
    elif ext in (".html", ".htm"):
        metadata = _extract_html_metadata(file_path)

    # Clean NUL bytes from all string fields to prevent DB errors
    for key in ["title", "author", "description"]:
        if key in metadata and isinstance(metadata[key], str):
            metadata[key] = metadata[key].replace("\x00", "")

    return metadata


def _extract_pptx_metadata(file_path: Path) -> dict:
    """Extract metadata from PPTX file using python-pptx."""
    metadata = {}
    try:
        from pptx import Presentation

        prs = Presentation(str(file_path))

        # Core properties
        core = prs.core_properties
        if core.title:
            metadata["title"] = core.title
        if core.author:
            metadata["author"] = core.author
        if core.comments:
            metadata["description"] = core.comments
        if core.created:
            metadata["created_date"] = core.created
        if core.modified:
            metadata["modified_date"] = core.modified

        # Slide count
        metadata["slide_count"] = len(prs.slides)

        # Try to get title from first slide
        if not metadata.get("title") and prs.slides:
            first_slide = prs.slides[0]
            for shape in first_slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text and len(text) > 3 and len(text) < 200:
                        metadata["title"] = text
                        break

    except Exception as e:
        logger.warning(f"Failed to extract PPTX metadata from {file_path.name}: {e}")

    return metadata


def _extract_pdf_metadata(file_path: Path) -> dict:
    """Extract metadata from PDF file using PyPDF2."""
    metadata = {}
    try:
        from PyPDF2 import PdfReader

        reader = PdfReader(str(file_path))

        # PDF metadata
        info = reader.metadata
        if info:
            if info.title:
                metadata["title"] = info.title
            if info.author:
                metadata["author"] = info.author
            if info.subject:
                metadata["description"] = info.subject
            if info.creation_date:
                metadata["created_date"] = info.creation_date

        # Page count
        metadata["slide_count"] = len(reader.pages)

    except Exception as e:
        logger.warning(f"Failed to extract PDF metadata from {file_path.name}: {e}")

    return metadata


def _extract_html_metadata(file_path: Path) -> dict:
    """Extract metadata from HTML presentation."""
    metadata = {}
    try:
        content = file_path.read_text(encoding="utf-8", errors="ignore")

        # Extract title from <title> tag
        import re
        title_match = re.search(r'<title[^>]*>(.*?)</title>', content, re.IGNORECASE | re.DOTALL)
        if title_match:
            metadata["title"] = title_match.group(1).strip()

        # Extract meta description
        desc_match = re.search(r'<meta\s+name=["\']description["\']\s+content=["\'](.*?)["\']', content, re.IGNORECASE)
        if desc_match:
            metadata["description"] = desc_match.group(1).strip()

        # Extract meta author
        author_match = re.search(r'<meta\s+name=["\']author["\']\s+content=["\'](.*?)["\']', content, re.IGNORECASE)
        if author_match:
            metadata["author"] = author_match.group(1).strip()

        # Estimate slide count from common slide markers
        slide_markers = re.findall(r'class=["\'][^"\']*slide[^"\']*["\']', content, re.IGNORECASE)
        if slide_markers:
            metadata["slide_count"] = len(slide_markers)

    except Exception as e:
        logger.warning(f"Failed to extract HTML metadata from {file_path.name}: {e}")

    return metadata
