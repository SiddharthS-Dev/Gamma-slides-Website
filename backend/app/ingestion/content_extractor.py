"""Deep content extraction engine — extracts full text, headings, and speaker notes.

Supports PPTX, PDF, and HTML files. Returns structured content for AI analysis.
This goes beyond the basic metadata_extractor by extracting full slide text,
speaker notes, headings, and body content separately for optimal classification.
"""

import logging
import re
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)


@dataclass
class ExtractedContent:
    """Structured content extracted from a presentation file."""
    title: str = ""
    headings: list[str] = field(default_factory=list)
    body_text: str = ""
    speaker_notes: str = ""
    metadata: dict = field(default_factory=dict)
    slide_texts: list[str] = field(default_factory=list)
    word_count: int = 0
    language: str = "en"
    file_name: str = ""
    file_type: str = ""


def extract_full_content(file_path: Path) -> ExtractedContent:
    """Extract deep content from a presentation file.

    This is the primary entry point for the content extraction pipeline.
    Unlike the basic metadata extractor, this extracts ALL text content
    for AI classification, tagging, and summarization.
    """
    ext = file_path.suffix.lower()
    content = ExtractedContent(
        file_name=file_path.name,
        file_type=ext.lstrip('.'),
    )

    try:
        if ext == ".pptx":
            content = _extract_pptx_content(file_path, content)
        elif ext == ".pdf":
            content = _extract_pdf_content(file_path, content)
        elif ext in (".html", ".htm"):
            content = _extract_html_content(file_path, content)
        else:
            logger.warning(f"Unsupported file type for content extraction: {ext}")

        # Compute word count
        all_text = " ".join([content.body_text or "", content.speaker_notes or ""])
        content.word_count = len(all_text.split())

        # Clean NUL bytes from all string fields to prevent DB errors
        if content.title:
            content.title = content.title.replace("\x00", "")
        if content.body_text:
            content.body_text = content.body_text.replace("\x00", "")
        if content.speaker_notes:
            content.speaker_notes = content.speaker_notes.replace("\x00", "")
        if content.headings:
            content.headings = [h.replace("\x00", "") for h in content.headings if h]
        if content.slide_texts:
            content.slide_texts = [s.replace("\x00", "") for s in content.slide_texts if s]
        if "description" in content.metadata and isinstance(content.metadata["description"], str):
            content.metadata["description"] = content.metadata["description"].replace("\x00", "")
        if "author" in content.metadata and isinstance(content.metadata["author"], str):
            content.metadata["author"] = content.metadata["author"].replace("\x00", "")

    except Exception as e:
        logger.error(f"Content extraction failed for {file_path.name}: {e}")

    return content


def _extract_pptx_content(file_path: Path, content: ExtractedContent) -> ExtractedContent:
    """Extract full content from PPTX including slide text and speaker notes."""
    try:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation(str(file_path))

        # Core metadata
        core = prs.core_properties
        if core.title:
            content.title = core.title
        if core.author:
            content.metadata["author"] = core.author
        if core.comments:
            content.metadata["description"] = core.comments
        if core.created:
            content.metadata["created_date"] = str(core.created)
        if core.modified:
            content.metadata["modified_date"] = str(core.modified)
        content.metadata["slide_count"] = len(prs.slides)

        all_body_parts = []
        all_notes_parts = []
        all_headings = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_text_parts = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # Detect headings: title shapes, or large font text
                is_heading = False

                # Check if this is a title placeholder
                if getattr(shape, 'is_placeholder', False):
                    try:
                        ph_type = shape.placeholder_format.type
                        # Placeholder types 0 (TITLE), 1 (CENTER_TITLE), 13 (TITLE)
                        if ph_type is not None and ph_type in (0, 1, 13, 15):
                            is_heading = True
                    except Exception:
                        pass

                # Also check by font size — anything >= 20pt is likely a heading
                if not is_heading:
                    try:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.size and run.font.size >= Pt(20):
                                    is_heading = True
                                    break
                            if is_heading:
                                break
                    except Exception:
                        pass

                if is_heading and text and len(text) < 200:
                    all_headings.append(text)

                slide_text_parts.append(text)

            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    all_notes_parts.append(notes_text)

            slide_text = " ".join(slide_text_parts)
            content.slide_texts.append(slide_text)
            all_body_parts.append(slide_text)

        # Set title from first heading if not from metadata
        if not content.title and all_headings:
            content.title = all_headings[0]

        content.headings = all_headings
        content.body_text = "\n\n".join(all_body_parts)
        content.speaker_notes = "\n\n".join(all_notes_parts)

    except ImportError:
        logger.warning("python-pptx not installed, skipping PPTX content extraction")
    except Exception as e:
        logger.error(f"PPTX content extraction failed: {e}")

    return content


def _extract_pdf_content(file_path: Path, content: ExtractedContent) -> ExtractedContent:
    """Extract full content from PDF with heading detection."""
    try:
        from PyPDF2 import PdfReader

        reader = PdfReader(str(file_path))

        # Metadata
        info = reader.metadata
        if info:
            if info.title:
                content.title = info.title
            if info.author:
                content.metadata["author"] = info.author
            if info.subject:
                content.metadata["description"] = info.subject
        content.metadata["slide_count"] = len(reader.pages)

        all_body_parts = []

        for page_idx, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            page_text = page_text.strip()

            if page_text:
                content.slide_texts.append(page_text)
                all_body_parts.append(page_text)

                # Heuristic heading detection: first line of each page if it's short
                lines = page_text.split('\n')
                for line in lines[:3]:  # Check first 3 lines
                    line = line.strip()
                    if (line and len(line) < 150 and len(line) > 3
                            and not line[0].isdigit()
                            and not line.startswith(('•', '-', '*', '●'))):
                        # Check if it looks like a heading (mostly uppercase or title case)
                        word_count = len(line.split())
                        if word_count <= 10:
                            content.headings.append(line)
                            break

        # Title from first heading if not from metadata
        if not content.title and content.headings:
            content.title = content.headings[0]

        content.body_text = "\n\n".join(all_body_parts)

    except ImportError:
        logger.warning("PyPDF2 not installed, skipping PDF content extraction")
    except Exception as e:
        logger.error(f"PDF content extraction failed: {e}")

    return content


def _extract_html_content(file_path: Path, content: ExtractedContent) -> ExtractedContent:
    """Extract content from HTML presentation files."""
    try:
        raw_html = file_path.read_text(encoding="utf-8", errors="ignore")

        # Try BeautifulSoup first, fall back to regex
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(raw_html, "html.parser")

            # Title
            title_tag = soup.find("title")
            if title_tag and title_tag.string:
                content.title = title_tag.string.strip()

            # Meta tags
            meta_desc = soup.find("meta", attrs={"name": "description"})
            if meta_desc and meta_desc.get("content"):
                content.metadata["description"] = meta_desc["content"]

            meta_author = soup.find("meta", attrs={"name": "author"})
            if meta_author and meta_author.get("content"):
                content.metadata["author"] = meta_author["content"]

            # Extract headings
            for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
                heading_text = tag.get_text(strip=True)
                if heading_text and len(heading_text) > 2:
                    content.headings.append(heading_text)

            # Extract body text (paragraphs, lists, divs)
            body_parts = []
            for tag in soup.find_all(["p", "li", "td", "th", "dd"]):
                text = tag.get_text(strip=True)
                if text and len(text) > 5:
                    body_parts.append(text)

            content.body_text = "\n".join(body_parts)

            # Detect slide sections
            slide_elements = soup.find_all(class_=re.compile(r'slide', re.IGNORECASE))
            if slide_elements:
                content.metadata["slide_count"] = len(slide_elements)
                for slide_el in slide_elements:
                    content.slide_texts.append(slide_el.get_text(strip=True))

        except ImportError:
            # Fallback to regex-based extraction
            logger.info("BeautifulSoup not available, using regex HTML extraction")
            content = _extract_html_regex(raw_html, content)

    except Exception as e:
        logger.error(f"HTML content extraction failed: {e}")

    return content


def _extract_html_regex(html: str, content: ExtractedContent) -> ExtractedContent:
    """Regex-based HTML extraction fallback when BeautifulSoup is not available."""
    # Title
    title_match = re.search(r'<title[^>]*>(.*?)</title>', html, re.IGNORECASE | re.DOTALL)
    if title_match:
        content.title = re.sub(r'<[^>]+>', '', title_match.group(1)).strip()

    # Headings
    heading_matches = re.findall(r'<h[1-6][^>]*>(.*?)</h[1-6]>', html, re.IGNORECASE | re.DOTALL)
    for h in heading_matches:
        text = re.sub(r'<[^>]+>', '', h).strip()
        if text and len(text) > 2:
            content.headings.append(text)

    # Body text — strip all tags
    body = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.IGNORECASE | re.DOTALL)
    body = re.sub(r'<style[^>]*>.*?</style>', '', body, flags=re.IGNORECASE | re.DOTALL)
    body = re.sub(r'<[^>]+>', ' ', body)
    body = re.sub(r'\s+', ' ', body).strip()
    content.body_text = body

    # Meta description
    desc_match = re.search(
        r'<meta\s+name=["\']description["\']\s+content=["\'](.*?)["\']',
        html, re.IGNORECASE,
    )
    if desc_match:
        content.metadata["description"] = desc_match.group(1).strip()

    # Slide count from slide markers
    slide_markers = re.findall(r'class=["\'][^"\']*slide[^"\']*["\']', html, re.IGNORECASE)
    if slide_markers:
        content.metadata["slide_count"] = len(slide_markers)

    return content
